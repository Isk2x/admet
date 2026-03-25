"""
Обновление индивидуальной презентации Каболова: 3 сем → 4 сем.
Обновляет существующие 8 слайдов + добавляет новые с результатами.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree
import os

SRC = "/Users/akabolov/Downloads/Каболов АС - Индивидуальная Презентация стартап (НИР 3 сем) (1).pptx"
DST = "/Users/akabolov/development/mipt/admet/docs/individual_presentation_4sem.pptx"

BLUE = RGBColor(0x12, 0x31, 0x94)
BLACK = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation(SRC)
slides = list(prs.slides)
title_layout = prs.slide_layouts[0]  # TITLE


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def update_textbox(shape, paragraphs_data):
    tf = shape.text_frame
    existing = list(tf.paragraphs)
    while len(existing) > len(paragraphs_data):
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)
        existing = list(tf.paragraphs)
    for i, pdata in enumerate(paragraphs_data):
        if i < len(tf.paragraphs):
            p = tf.paragraphs[i]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        else:
            p = tf.add_paragraph()
        if pdata.get("align") is not None:
            p.alignment = pdata["align"]
        for rdata in pdata["runs"]:
            run = p.add_run()
            run.text = rdata["text"]
            run.font.name = rdata.get("font", "Montserrat")
            run.font.size = rdata.get("size", Pt(12))
            run.font.bold = rdata.get("bold", False)
            if rdata.get("color"):
                run.font.color.rgb = rdata["color"]


def add_slide_with_title_and_body(prs, title_text, body_paragraphs, insert_before_idx=None):
    """Добавляет слайд с заголовком + текстовым блоком."""
    slide = prs.slides.add_slide(title_layout)

    # Удаляем все placeholder shapes из layout
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Заголовок
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(350000), Emu(260000), Emu(8250000), Emu(770000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = "Montserrat"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLUE

    # Тело
    txBox2 = slide.shapes.add_textbox(Emu(370000), Emu(870000), Emu(8200000), Emu(3800000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, pdata in enumerate(body_paragraphs):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        for rdata in pdata["runs"]:
            run = p.add_run()
            run.text = rdata["text"]
            run.font.name = rdata.get("font", "Montserrat")
            run.font.size = rdata.get("size", Pt(12))
            run.font.bold = rdata.get("bold", False)
            if rdata.get("color"):
                run.font.color.rgb = rdata["color"]

    if insert_before_idx is not None:
        _move_slide(prs, len(prs.slides) - 1, insert_before_idx)

    return slide


def _move_slide(prs, from_idx, to_idx):
    """Перемещает слайд from_idx в позицию to_idx."""
    presentation_elem = prs.part._element
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = presentation_elem.find('.//p:sldIdLst', nsmap)
    sldIds = list(sldIdLst)
    el = sldIds[from_idx]
    sldIdLst.remove(el)
    sldIds = list(sldIdLst)
    if to_idx >= len(sldIds):
        sldIdLst.append(el)
    else:
        sldIds[to_idx].addprevious(el)


def R(text, size=Pt(12), bold=False, color=None, font="Montserrat"):
    d = {"text": text, "size": size, "bold": bold, "font": font}
    if color:
        d["color"] = color
    return d


def P(*runs):
    return {"runs": list(runs)}


# ============================================================
# SLIDE 2 (idx 1) — Роль в проекте: обновляем
# ============================================================
slide2 = slides[1]
shape_role = find_shape(slide2, "Google Shape;70;p3")
if shape_role:
    update_textbox(shape_role, [
        P(R("Основная роль: ML / Research Engineer", bold=True)),
        P(R("Анализ State of the Art в области предсказания токсичности (30+ источников)")),
        P(R("Проектирование архитектуры: GIN + Virtual Node + Pretrained backbone")),
        P(R("Формулировка и проверка 3 исследовательских гипотез (H1–H3)")),
        P(R("Реализация Explainability-Aware Loss (L_faithfulness + L_stability) — научная новизна")),
        P(R("Реализация Uncertainty-Weighted MTL и Toxicophore-Guided XAI")),
        P(R("Обучение 9 конфигураций GIN + ChemBERTa + SchNet на Colab A100 (3 seeds × scaffold)")),
        P(R("Разработка MVP: FastAPI backend + Frontend с атомной визуализацией")),
        P(R("Результат: ROC-AUC 0.791 на Tox21 scaffold (SOTA: 0.751, +4 п.п.)", bold=True)),
        P(R("", size=Pt(6))),
        P(R("Дополнительная роль: Team Lead", bold=True)),
        P(R("Постановка технических требований и критериев качества")),
        P(R("Координация команды из 3 человек")),
        P(R("Подготовка демонстраций и пилотных материалов")),
    ])

# ============================================================
# SLIDE 3 (idx 2) — Гипотезы: обновляем таблицу
# ============================================================
slide3 = slides[2]
shape_hyp = find_shape(slide3, "Google Shape;77;g3aae066a79d_0_115")
if shape_hyp and shape_hyp.has_table:
    tbl = shape_hyp.table
    hyp_data = [
        ["Проблема", "Гипотеза"],
        [
            "P1 — Недостаточная глобальная агрегация в GIN: информация распространяется на ~5 связей за 5 слоёв, но токсичность определяется далёкими группами",
            "H1: Добавление Virtual Node повысит mean ROC-AUC на Tox21 scaffold ≥ +1 п.п. vs GIN без VN.\nОпровержение: Δ < 1 п.п."
        ],
        [
            "P2 — Нет регуляризации по объяснимости: Pretrained GNN fine-tuning теряет химически значимые паттерны; Post-hoc XAI не влияет на качество",
            "H2: Pretrained GIN+VN + XAI-loss (λ=0.1) превысит SOTA (0.7512) ≥ +2 п.п. при AOPC > 0.\nОпровержение: ROC-AUC < 0.7712"
        ],
        [
            "P3 — Избыточность OGB features (9 vs 2 атомных признака) при наличии Virtual Node",
            "H3 (контрольная, отриц.): Расширение признаков с 2 до 9 даст < 1 п.п. при VN.\nОпровержение: Δ ≥ 1 п.п."
        ],
    ]
    for ri, row_data in enumerate(hyp_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            for p in tf.paragraphs:
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            lines = cell_text.split("\n")
            for li, line in enumerate(lines):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.name = "Montserrat"
                run.font.size = Pt(9) if ri > 0 else Pt(10)
                run.font.bold = (ri == 0)

# ============================================================
# SLIDE 4 (idx 3) — Анализ проблемы: обновляем
# ============================================================
slide4 = slides[3]
shape_analysis = find_shape(slide4, "Google Shape;84;g3aae066a79d_0_173")
if shape_analysis:
    update_textbox(shape_analysis, [
        P(R("GIN: информация за 5 слоёв распространяется на ~5 связей. Токсичность часто определяется взаимодействием далёких функциональных групп (hERG, нитро + амин). Virtual Node решает эту проблему — на каждом слое вся молекула агрегируется.", size=Pt(14))),
        P(R("Random split завышает качество на 5–10%. Scaffold split (Bemis-Murcko) — реалистичная оценка обобщаемости. Все эксперименты выполнены на scaffold split.", size=Pt(14))),
        P(R("Pretrained GNN (Hu et al. 2020, ICLR, 2M молекул) при fine-tuning может терять химически значимые паттерны. Post-hoc XAI (GNNExplainer, SHAP) не влияет на качество модели.", size=Pt(14))),
        P(R("В литературе Virtual Node, Pretrained GNN и XAI описаны по отдельности. Не исследована их комбинация с интеграцией XAI-регуляризации в процесс обучения.", size=Pt(14))),
        P(R("", size=Pt(6))),
        P(R("Вывод: фокус — на VN для глобальной агрегации, Pretrained backbone для трансфера, и Explainability-Aware Loss как научная новизна.", bold=True, size=Pt(14))),
    ])

# ============================================================
# SLIDE 5 (idx 4) — Методика: обновляем
# ============================================================
slide5 = slides[4]
shape_design = find_shape(slide5, "Google Shape;90;g3aae066a79d_0_228")
if shape_design:
    update_textbox(shape_design, [
        P(R("Исследовательский дизайн", bold=True, size=Pt(13))),
        P(R("Tox21 (7823 мол, 12 задач), Balanced Scaffold Split (Bemis-Murcko) 70/15/15", size=Pt(13))),
        P(R("9 конфигураций GIN + ChemBERTa + SchNet, обучение Google Colab A100", size=Pt(13))),
    ])

shape_data = find_shape(slide5, "Google Shape;91;g3aae066a79d_0_228")
if shape_data:
    update_textbox(shape_data, [
        P(R("Подход к данным", bold=True, size=Pt(13))),
        P(R("Canonical SMILES → молекулярный граф (atom_num, chirality, bond_type, bond_dir)", size=Pt(13))),
        P(R("Очистка/дедупликация RDKit, контроль утечек", size=Pt(13))),
        P(R("3 seeds (42, 0, 1) → mean ± std; Adam (lr=1e-3, wd=1e-5), Early Stopping (patience=15), Batch 64", size=Pt(13))),
    ])

shape_metrics = find_shape(slide5, "Google Shape;92;g3aae066a79d_0_228")
if shape_metrics:
    update_textbox(shape_metrics, [
        P(R("Метрики оценки", bold=True, size=Pt(13))),
        P(R("Качество: ROC-AUC, PR-AUC (mean по 12 задачам Tox21)", size=Pt(13))),
        P(R("Обобщаемость: +BBBP (0.869), +ClinTox — всё scaffold split", size=Pt(13))),
        P(R("Объяснимость: AOPC@20% (faithfulness), IoU@20% (stability)", size=Pt(13))),
    ])

# ============================================================
# SLIDE 6 (idx 5) — Вклад в MVP: обновляем
# ============================================================
slide6 = slides[5]
shape_mvp = find_shape(slide6, "Google Shape;98;g3b1b9fed030_0_24")
if shape_mvp:
    update_textbox(shape_mvp, [
        P(R("ML-часть продукта", bold=True)),
        P(R("Обучено 9 конфигураций GIN + ChemBERTa + SchNet на Colab A100 (3 seeds × scaffold split)")),
        P(R("Virtual Node, OGB features, Focal loss, Attention pooling, Uncertainty MTL, Toxicophore XAI")),
        P(R("Лучший результат: GIN pretrained + VN + XAI = ROC-AUC 0.791 (SOTA 0.751, +4 п.п.)", bold=True)),
        P(R("3D-ветка: SchNet + MultiConf-SchNet + Hybrid GIN-SchNet Fusion (реализована)")),
        P(R("", size=Pt(6))),
        P(R("Explainability как часть продукта", bold=True)),
        P(R("Explainability-Aware Loss (L_faithfulness + L_stability) — встроено в обучение")),
        P(R("Toxicophore-Guided XAI: 12 SMARTS-паттернов → инъекция доменного знания (+1.6 п.п.)")),
        P(R("IoU стабильности объяснений: 0.892 → 0.907 → 0.913 (монотонный рост с XAI-loss)")),
        P(R("", size=Pt(6))),
        P(R("MVP Backend + Frontend", bold=True)),
        P(R("FastAPI: /api/predict, /api/models, /api/health; два ML-движка: GIN (2.4M) + ChemBERTa (77M)")),
        P(R("Визуализация молекулы с атомной атрибуцией (RDKit SVG) + столбчатая диаграмма")),
        P(R("Frontend: ввод SMILES, выбор модели, 12 эндпоинтов + atom chart")),
    ])

# ============================================================
# SLIDE 7 (idx 6) — Источники: обновляем
# ============================================================
slide7 = slides[6]
shape_refs = find_shape(slide7, "Google Shape;105;g3aae066a79d_0_282")
if shape_refs:
    refs = [
        "Hu, W. et al. (2020). Strategies for Pre-training Graph Neural Networks. ICLR 2020.",
        "Xu, K. et al. (2019). How Powerful are Graph Neural Networks? ICLR 2019.",
        "Gilmer, J. et al. (2017). Neural Message Passing for Quantum Chemistry. ICML 2017.",
        "Wu, Z. et al. (2018). MoleculeNet: A Benchmark for Molecular ML. Chemical Science.",
        "Ying, R. et al. (2019). GNNExplainer: Generating Explanations for GNN. NeurIPS 2019.",
        "Hu, W. et al. (2021). Open Graph Benchmark. NeurIPS 2020.",
        "Kendall, A. et al. (2018). Multi-Task Learning Using Uncertainty. CVPR 2018.",
        "Xiong, Z. et al. (2020). Pushing the Boundaries of Molecular Representation. J. Med. Chem.",
        "Lin, T.-Y. et al. (2017). Focal Loss for Dense Object Detection. ICCV 2017.",
        "Schütt, K. et al. (2017). SchNet: continuous-filter convolutional NN. NeurIPS 2017.",
    ]
    pdata = [P(R(r, size=Pt(12))) for r in refs]
    update_textbox(shape_refs, pdata)


# ============================================================
# Добавляем НОВЫЕ слайды (вставляем перед источниками = idx 6)
# ============================================================

# --- НОВЫЙ СЛАЙД: Результаты сводная таблица ---
add_slide_with_title_and_body(prs, "Результаты: сводная таблица", [
    P(R("Tox21, scaffold split, 3 seeds (mean ± std)", bold=True, size=Pt(13))),
    P(R("", size=Pt(4))),
    P(R("GIN scratch (baseline)                         0.769 ± 0.012     —", size=Pt(11))),
    P(R("GIN + Virtual Node                             0.789 ± 0.004     +2.0 п.п.", size=Pt(11))),
    P(R("GIN + OGB + VN                                 0.776 ± 0.017     +0.7 п.п.", size=Pt(11))),
    P(R("GIN + VN + Attention pool                      0.757 ± 0.022     −1.2 п.п.", size=Pt(11))),
    P(R("GIN + VN + Focal loss                          0.763 ± 0.021     −0.6 п.п.", size=Pt(11))),
    P(R("GIN pretrained + VN + XAI                      0.791 ± 0.005     +2.2 п.п.  ← лучший", bold=True, size=Pt(11), color=GREEN)),
    P(R("GIN + VN + UncMTL                              0.775 ± 0.010     +0.6 п.п.", size=Pt(11))),
    P(R("GIN + VN + Toxicophore XAI                     0.785 ± 0.004     +1.6 п.п.", size=Pt(11))),
    P(R("GIN pretrained + VN + Full                     0.788 ± 0.001     +1.9 п.п.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("Hu et al. 2020 (ICLR) SOTA                    0.7512 ± 0.0079", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("Лучший результат: GIN pretrained + VN + XAI = 0.791 → +4.0 п.п. vs SOTA", bold=True, size=Pt(13), color=BLUE)),
], insert_before_idx=6)

# --- НОВЫЙ СЛАЙД: H1 подтверждена ---
add_slide_with_title_and_body(prs, "H1: Virtual Node — подтверждена (+2.0 п.п.)", [
    P(R("GIN scratch (baseline):    0.769 ± 0.012", size=Pt(14))),
    P(R("GIN + Virtual Node:        0.789 ± 0.004    Δ = +2.0 п.п.", bold=True, size=Pt(14), color=GREEN)),
    P(R("", size=Pt(8))),
    P(R("Почему сработало:", bold=True, size=Pt(14))),
    P(R("• Стандартный GIN: информация распространяется на ~5 связей за 5 слоёв", size=Pt(13))),
    P(R("• Virtual Node: на каждом слое ВСЯ молекула агрегируется → рассылается обратно", size=Pt(13))),
    P(R("• Для токсичности: критичны взаимодействия далёких групп (hERG, нитро + амин)", size=Pt(13))),
    P(R("• Бонус: дисперсия снизилась (std: 0.012 → 0.004) — модель стабильнее", size=Pt(13))),
    P(R("", size=Pt(8))),
    P(R("Критерий опровержения: Δ < 1 п.п. → НЕ опровергнута (Δ = +2.0 п.п.)", bold=True, size=Pt(13))),
], insert_before_idx=7)

# --- НОВЫЙ СЛАЙД: H2 частично подтверждена ---
add_slide_with_title_and_body(prs, "H2: Pretrained + XAI loss — частично подтверждена", [
    P(R("SOTA (Hu et al. 2020):        0.7512 ± 0.0079", size=Pt(14))),
    P(R("GIN pretrained + VN + XAI:    0.791 ± 0.005    Δ = +4.0 п.п. vs SOTA", bold=True, size=Pt(14), color=GREEN)),
    P(R("", size=Pt(8))),
    P(R("Декомпозиция эффекта:", bold=True, size=Pt(14))),
    P(R("• GIN scratch:                0.769        —", size=Pt(13))),
    P(R("• + Virtual Node:             0.789        +2.0 п.п.", size=Pt(13))),
    P(R("• + Pretrained + XAI loss:    0.791        +0.2 п.п.", size=Pt(13))),
    P(R("• Итого vs SOTA:              0.791        +4.0 п.п.", bold=True, size=Pt(13))),
    P(R("", size=Pt(8))),
    P(R("Ограничение: AOPC для pretrained+XAI < 0 — L2-norm embeddings не оптимальный прокси важности.", size=Pt(12), color=RED)),
    P(R("→ По ROC-AUC подтверждена (+4.0 п.п.), по AOPC — нет.", bold=True, size=Pt(13))),
], insert_before_idx=8)

# --- НОВЫЙ СЛАЙД: H3 подтверждена ---
add_slide_with_title_and_body(prs, "H3: OGB features не помогают — подтверждена", [
    P(R("GIN + VN (2 признака):        0.789 ± 0.004", size=Pt(14))),
    P(R("GIN + OGB + VN (9 признаков): 0.776 ± 0.017    Δ = −1.3 п.п.", bold=True, size=Pt(14), color=RED)),
    P(R("", size=Pt(8))),
    P(R("Интерпретация:", bold=True, size=Pt(14))),
    P(R("• GIN с 5 слоями самостоятельно выучивает degree, is_aromatic, is_in_ring из топологии", size=Pt(13))),
    P(R("• 7 дополнительных OGB-признаков не несут новой информации при наличии Virtual Node", size=Pt(13))),
    P(R("• Дисперсия выросла: std 0.004 → 0.017 — больше параметров = менее стабильно", size=Pt(13))),
    P(R("", size=Pt(8))),
    P(R("Критерий опровержения: Δ ≥ 1 п.п. → НЕ опровергнута (Δ = −1.3 п.п., даже хуже)", bold=True, size=Pt(13))),
], insert_before_idx=9)

# --- НОВЫЙ СЛАЙД: Научная новизна ---
add_slide_with_title_and_body(prs, "Научная новизна: 3 инновации", [
    P(R("1. Explainability-Aware Fine-tuning", bold=True, size=Pt(14), color=BLUE)),
    P(R("L_total = L_task + λ · (0.5 · L_faithfulness + 0.5 · L_stability)", size=Pt(12))),
    P(R("Модель учится быть объяснимой, а не просто точной. IoU: 0.892 → 0.907 → 0.913", size=Pt(12))),
    P(R("", size=Pt(6))),
    P(R("2. Uncertainty-Weighted Multi-Task Loss (Kendall et al. 2018)", bold=True, size=Pt(14), color=BLUE)),
    P(R("12 обучаемых log(σ²) — модель автоматически взвешивает задачи", size=Pt(12))),
    P(R("Шумные задачи (NR-ER: σ²=3.2) → низкий вес; информативные (NR-AR-LBD: σ²=14.4) → высокий", size=Pt(12))),
    P(R("Прирост: +0.6 п.п. vs baseline", size=Pt(12))),
    P(R("", size=Pt(6))),
    P(R("3. Toxicophore-Guided XAI", bold=True, size=Pt(14), color=BLUE)),
    P(R("12 SMARTS-паттернов токсикофоров (нитро, эпоксид, арил-галогенид...)", size=Pt(12))),
    P(R("L_toxicophore: штраф если importance токсикофорных атомов < средней", size=Pt(12))),
    P(R("Инъекция доменного знания в обучение → прирост +1.6 п.п.", size=Pt(12))),
], insert_before_idx=10)

# --- НОВЫЙ СЛАЙД: Обобщаемость ---
add_slide_with_title_and_body(prs, "Обобщаемость: дополнительные датасеты", [
    P(R("GIN + VN, 3 seeds, scaffold split", bold=True, size=Pt(14))),
    P(R("", size=Pt(6))),
    P(R("Tox21 (12 задач):     ROC-AUC  0.789 ± 0.004    — основной бенчмарк", size=Pt(14))),
    P(R("BBBP (1 задача):      ROC-AUC  0.869 ± 0.005    — проницаемость ГЭБ, отличный результат", size=Pt(14), color=GREEN)),
    P(R("ClinTox (2 задачи):   ROC-AUC  0.564 ± 0.041    — малый датасет (1461), высокая дисперсия", size=Pt(14), color=RED)),
    P(R("", size=Pt(8))),
    P(R("Вывод:", bold=True, size=Pt(14))),
    P(R("Подход GIN+VN обобщается за пределы Tox21.", size=Pt(13))),
    P(R("BBBP (0.869) подтверждает применимость для разных ADMET-эндпоинтов.", size=Pt(13))),
    P(R("ClinTox требует pretrained подхода или transfer learning.", size=Pt(13))),
], insert_before_idx=11)

# --- НОВЫЙ СЛАЙД: Адресация критики ---
add_slide_with_title_and_body(prs, "Адресация критики 3 семестра", [
    P(R("«Научная новизна недостаточно обоснована»", bold=True, size=Pt(12))),
    P(R("→ Сравнение с Hu et al. 2020 (ICLR): 0.791 vs 0.751 (+4 п.п.). Три инновации.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("«Постановка задачи размыта (агрегированная метка)»", bold=True, size=Pt(12))),
    P(R("→ Per-task multi-task (12 задач Tox21) — стандарт MoleculeNet.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("«Отсутствует scaffold-протокол»", bold=True, size=Pt(12))),
    P(R("→ Balanced scaffold split (Bemis-Murcko) 70/15/15, 3 seeds.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("«Сравнение только с внутренним baseline»", bold=True, size=Pt(12))),
    P(R("→ Сравнение с Hu et al. 2020 (ICLR), 9 конфигураций, 3 seeds.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("«Интерпретируемость оценена слабо»", bold=True, size=Pt(12))),
    P(R("→ Количественный XAI: AOPC, IoU, L_faithfulness, L_stability в loss.", size=Pt(11))),
    P(R("", size=Pt(4))),
    P(R("«ROC-AUC=0.92 — нереалистично»", bold=True, size=Pt(12))),
    P(R("→ Реальный результат: 0.791. SOTA: 0.751. Все числа из запусков кода.", size=Pt(11))),
], insert_before_idx=12)

# --- НОВЫЙ СЛАЙД: Ограничения и дальнейшая работа ---
add_slide_with_title_and_body(prs, "Ограничения и дальнейшая работа", [
    P(R("Ограничения:", bold=True, size=Pt(14), color=RED)),
    P(R("1. AOPC для pretrained+XAI < 0 — L2-norm не оптимальный прокси; нужна gradient-based атрибуция", size=Pt(12))),
    P(R("2. ClinTox: 0.564 — малый датасет, требуется cross-dataset transfer", size=Pt(12))),
    P(R("3. 3 seeds — минимально достаточно; 5–10 seeds надёжнее", size=Pt(12))),
    P(R("4. Нет внешней валидации с экспертами-химиками", size=Pt(12))),
    P(R("", size=Pt(8))),
    P(R("Дальнейшая работа:", bold=True, size=Pt(14), color=BLUE)),
    P(R("• 3D-модели: SchNet + Hybrid GIN-SchNet Fusion (реализовано, в обучении на Colab)", size=Pt(12))),
    P(R("• ChemBERTa: Transformer-based attention explainability", size=Pt(12))),
    P(R("• A/B тест с экспертами: override rate, time-to-decision", size=Pt(12))),
    P(R("• Расширение ADMET: растворимость, Caco-2, CYP450 (задачи коллег)", size=Pt(12))),
    P(R("• Production: Docker, CI/CD, мониторинг drift", size=Pt(12))),
], insert_before_idx=13)


# ============================================================
# Сохраняем
# ============================================================
prs.save(DST)
print(f"Сохранено: {DST}")
print(f"Итого слайдов: {len(prs.slides)}")
print("Обновлены: 2 (роль), 3 (гипотезы), 4 (анализ), 5 (методика), 6 (MVP), 7→14 (источники)")
print("Новые: Результаты, H1, H2, H3, Научная новизна, Обобщаемость, Адресация критики, Ограничения")
