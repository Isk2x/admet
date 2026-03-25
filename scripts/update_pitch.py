"""
Скрипт обновления питч-дека: 3 сем → 4 сем.
Обновляет общие слайды (1-7) и слайды Каболова (8-13).
Слайды остальных участников не трогаются.
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

SRC = "/Users/akabolov/Downloads/pitch_4sem.pptx"
DST = "/Users/akabolov/development/mipt/admet/docs/pitch_deck_4sem.pptx"

prs = Presentation(SRC)
slides = list(prs.slides)


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def set_text_simple(shape, lines, font_name="Montserrat", font_size=Pt(13), bold=False, color=None):
    """Заменяет весь текст в shape новыми строками."""
    tf = shape.text_frame
    for i, p in enumerate(tf.paragraphs):
        for r in list(p.runs):
            r.text = ""
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, str):
            parts = [(line_data, bold)]
        elif isinstance(line_data, list):
            parts = line_data
        else:
            parts = [(str(line_data), bold)]

        p.alignment = tf.paragraphs[0].alignment if i > 0 else p.alignment

        for j, (text, is_bold) in enumerate(parts):
            if j == 0 and p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = is_bold
            if color:
                run.font.color.rgb = color


def update_textbox(shape, paragraphs_data):
    """
    paragraphs_data: list of dicts:
      { "runs": [{"text": ..., "bold": True/False, "size": Pt(...), "font": "...", "color": RGBColor/None}],
        "align": alignment or None }
    """
    tf = shape.text_frame
    existing = list(tf.paragraphs)

    while len(existing) > len(paragraphs_data):
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)
        existing = list(tf.paragraphs)

    for i, pdata in enumerate(paragraphs_data):
        if i < len(tf.paragraphs):
            p = tf.paragraphs[i]
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
        else:
            p = tf.add_paragraph()

        if pdata.get("align") is not None:
            p.alignment = pdata["align"]

        for rdata in pdata["runs"]:
            run = p.add_run()
            run.text = rdata["text"]
            run.font.name = rdata.get("font", "Montserrat")
            run.font.size = rdata.get("size", Pt(13))
            run.font.bold = rdata.get("bold", False)
            if rdata.get("color"):
                run.font.color.rgb = rdata["color"]


# ============================================================
# SLIDE 4 (idx 3) — Описание решения: обновляем результаты
# ============================================================
slide4 = slides[3]
shape_solution = find_shape(slide4, "Google Shape;96;p5")
if shape_solution:
    update_textbox(shape_solution, [
        {"runs": [{"text": "Use case:", "bold": True, "size": Pt(17)}]},
        {"runs": [
            {"text": "Био-аналитик загружает библиотеку 10k молекул → платформа фильтрует 95% потенциально проблемных кандидатов по ADMET-профилю и возвращает приоритетный список 500 молекул с прогнозом, ", "size": Pt(17)},
            {"text": "атомно-ориентированным объяснением", "bold": True, "size": Pt(17)},
            {"text": " и доверительным интервалом.", "size": Pt(17)},
        ]},
        {"runs": [{"text": "Ключевые результаты:", "bold": True, "size": Pt(15)}]},
        {"runs": [{"text": "Время оценки: дни вместо нескольких месяцев.", "size": Pt(15)}]},
        {"runs": [{"text": "Экономия: сокращение синтезов ≥30%.", "size": Pt(15)}]},
        {"runs": [{"text": "Прозрачность: прогноз + атомная атрибуция + визуализация молекулы.", "size": Pt(15)}]},
        {"runs": [{"text": "Точность: ROC-AUC 0.791 на Tox21 scaffold (SOTA: 0.751, +4 п.п.).", "bold": True, "size": Pt(15)}]},
    ])

# ============================================================
# SLIDE 5 (idx 4) — УЦП: обновляем с конкретными результатами
# ============================================================
slide5 = slides[4]
shape_uvp = find_shape(slide5, "Google Shape;103;g3458075d622_0_0")
if shape_uvp:
    update_textbox(shape_uvp, [
        {"runs": [{"text": "Российская ИИ-платформа для комплексного ADMET-моделирования лекарств", "size": Pt(13)}]},
        {"runs": [{"text": "", "size": Pt(8)}]},
        {"runs": [{"text": "Ценность: ", "bold": True, "size": Pt(13)}]},
        {"runs": [{"text": "Снижение стоимости разработки препаратов", "size": Pt(13)}]},
        {"runs": [{"text": "Ускорение процесса отбора кандидатов", "size": Pt(13)}]},
        {"runs": [{"text": "Интерпретируемость: видно, КАКИЕ атомы вызывают токсичность", "size": Pt(13)}]},
        {"runs": [{"text": "Универсальная платформа для множества ADMET-эндпоинтов", "size": Pt(13)}]},
        {"runs": [{"text": "", "size": Pt(8)}]},
        {"runs": [{"text": "Наше преимущество:", "bold": True, "size": Pt(13)}]},
        {"runs": [{"text": "Graph Neural Networks (GIN + Virtual Node + Pretrained) — ROC-AUC 0.791 (SOTA 0.751)", "size": Pt(13)}]},
        {"runs": [{"text": "Explainable AI — атомная атрибуция с IoU стабильностью 0.91", "size": Pt(13)}]},
        {"runs": [{"text": "Полный ADMET-профиль — токсичность, растворимость, проницаемость, метаболизм", "size": Pt(13)}]},
        {"runs": [{"text": "On-premise — конфиденциальность данных заказчика", "size": Pt(13)}]},
        {"runs": [{"text": "Импортонезависимость — нет санкционных рисков", "size": Pt(13)}]},
        {"runs": [{"text": "", "size": Pt(8)}]},
        {"runs": [
            {"text": "Пилот: ", "bold": True, "size": Pt(13)},
            {"text": "Фармацевтические компании и исследовательские лаборатории", "size": Pt(13)},
        ]},
    ])

# ============================================================
# SLIDE 7 (idx 6) — Прогресс: обновляем таблицу
# ============================================================
slide7 = slides[6]
shape_progress = find_shape(slide7, "Google Shape;117;g38c0d71ece1_0_468")
if shape_progress and shape_progress.has_table:
    tbl = shape_progress.table
    progress_data = [
        ["Готовность MVP\nРаботающий backend (FastAPI) + frontend,\nдва ML-движка (GIN + ChemBERTa), атомная атрибуция", "80 %"],
        ["Тестирование на пользователях\nROC-AUC 0.791 на Tox21 scaffold,\nBBBP: 0.869, 3 seeds × scaffold split", "75 %"],
        ["Тестирование продуктовой гипотезы\n3 гипотезы H1–H3 проверены,\nSOTA превышена на +4 п.п.", "90 %"],
        ["Go-to-market стратегия\nЦелевые сегменты, каналы,\n3 этапа GTM определены", "50 %"],
        ["Финансовые модели\nTAM/SAM/SOM рассчитан,\nюнит-экономика и P&L", "40 %"],
    ]
    for ri, row_data in enumerate(progress_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            for p in tf.paragraphs:
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            lines = cell_text.split("\n")
            for li, line in enumerate(lines):
                if li == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.name = "Montserrat"
                if li == 0:
                    run.font.size = Pt(10)
                    run.font.bold = True
                else:
                    run.font.size = Pt(8)
                    run.font.bold = False

# ============================================================
# SLIDE 9 (idx 8) — Описание задач Каболова: обновляем
# ============================================================
slide9 = slides[8]
shape_role = find_shape(slide9, "Google Shape;131;g38c0d71ece1_0_131")
if shape_role:
    update_textbox(shape_role, [
        {"runs": [{"text": "Основная роль: ML-инженер (Toxicity & Explainability)", "bold": True, "size": Pt(12)}]},
        {"runs": [{"text": "Анализ SOTA (30+ источников), проектирование архитектуры GIN + Virtual Node + Pretrained backbone", "size": Pt(12)}]},
        {"runs": [{"text": "Формулировка и проверка 3 исследовательских гипотез (H1–H3)", "size": Pt(12)}]},
        {"runs": [{"text": "Реализация Explainability-Aware Loss (L_faithfulness + L_stability) — научная новизна", "size": Pt(12)}]},
        {"runs": [{"text": "Реализация Uncertainty-Weighted MTL и Toxicophore-Guided XAI", "size": Pt(12)}]},
        {"runs": [{"text": "Обучение 9 конфигураций GIN + ChemBERTa + SchNet на Google Colab A100 (3 seeds × scaffold)", "size": Pt(12)}]},
        {"runs": [{"text": "Разработка MVP: FastAPI backend + Frontend с атомной визуализацией", "size": Pt(12)}]},
        {"runs": [{"text": "Результат: ROC-AUC 0.791 на Tox21 scaffold (SOTA 0.751, +4 п.п.)", "bold": True, "size": Pt(12)}]},
        {"runs": [{"text": "", "size": Pt(6)}]},
        {"runs": [{"text": "Дополнительная роль: Project Lead", "bold": True, "size": Pt(12)}]},
        {"runs": [{"text": "Постановка технических требований и критериев качества", "size": Pt(12)}]},
        {"runs": [{"text": "Координация команды из 3 человек", "size": Pt(12)}]},
        {"runs": [{"text": "Подготовка демонстраций и пилотных материалов", "size": Pt(12)}]},
    ])

# ============================================================
# SLIDE 10 (idx 9) — Проблемы и гипотезы Каболова: обновляем таблицу
# ============================================================
slide10 = slides[9]
shape_hyp_table = find_shape(slide10, "Google Shape;140;g3bdf68daae8_0_0")
if shape_hyp_table and shape_hyp_table.has_table:
    tbl = shape_hyp_table.table
    hyp_data = [
        ["Проблема", "Гипотеза"],
        [
            "P1 — Недостаточная глобальная агрегация информации в GIN (5 слоёв = ~5 связей)",
            "H1: Добавление Virtual Node повысит ROC-AUC ≥ +1 п.п. vs GIN без VN"
        ],
        [
            "P2 — Нет регуляризации по объяснимости; post-hoc XAI не влияет на качество модели",
            "H2: Pretrained GIN+VN + XAI-loss (λ=0.1) превысит SOTA (0.7512) ≥ +2 п.п."
        ],
        [
            "P3 — Избыточность расширенных атомных признаков (OGB, 9 вместо 2) при наличии VN",
            "H3 (контрольная, отриц.): OGB features дадут прирост < 1 п.п. при VN"
        ],
    ]
    for ri, row_data in enumerate(hyp_data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            for p in tf.paragraphs:
                for r in list(p.runs):
                    r._r.getparent().remove(r._r)
            while len(tf.paragraphs) > 1:
                tf.paragraphs[-1]._p.getparent().remove(tf.paragraphs[-1]._p)
            run = tf.paragraphs[0].add_run()
            run.text = cell_text
            run.font.name = "Montserrat"
            run.font.size = Pt(9) if ri > 0 else Pt(10)
            run.font.bold = (ri == 0)

# ============================================================
# SLIDE 11 (idx 10) — Анализ проблемы: обновляем
# ============================================================
slide11 = slides[10]
shape_analysis = find_shape(slide11, "Google Shape;147;g3bdf68daae8_0_55")
if shape_analysis:
    update_textbox(shape_analysis, [
        {"runs": [{"text": "GIN: информация распространяется на ~5 связей за 5 слоёв. Токсичность часто определяется взаимодействием далёких функциональных групп (hERG, нитро + амин). Virtual Node решает эту проблему.", "size": Pt(14)}]},
        {"runs": [{"text": "Random split завышает качество на 5–10%. Scaffold split (Bemis-Murcko) — реалистичная оценка обобщаемости. Все наши эксперименты — на scaffold split.", "size": Pt(14)}]},
        {"runs": [{"text": "Pretrained GNN (Hu et al. 2020, 2M молекул) + fine-tuning может терять химически значимые паттерны. Post-hoc XAI (GNNExplainer, SHAP) не влияет на качество модели.", "size": Pt(14)}]},
        {"runs": [{"text": "В литературе Virtual Node, Pretrained GNN и XAI описаны по отдельности. Не исследована их комбинация с интеграцией XAI-регуляризации в процесс обучения.", "size": Pt(14)}]},
        {"runs": [{"text": "", "size": Pt(8)}]},
        {"runs": [
            {"text": "Вывод: фокус — на VN для глобальной агрегации, Pretrained backbone для трансфера знаний, и Explainability-Aware Loss как научная новизна.", "bold": True, "size": Pt(14)},
        ]},
    ])

# ============================================================
# SLIDE 12 (idx 11) — Методика: обновляем
# ============================================================
slide12 = slides[11]
# Блок "Исследовательский дизайн"
shape_design = find_shape(slide12, "Google Shape;153;g3bdf68daae8_0_110")
if shape_design:
    update_textbox(shape_design, [
        {"runs": [{"text": "Исследовательский дизайн", "bold": True, "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Единый ML-контур: Tox21 (7823 мол, 12 задач), scaffold split 70/15/15, 3 seeds (42, 0, 1)", "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "9 конфигураций GIN + ChemBERTa + SchNet, обучение на Google Colab A100", "size": Pt(13), "font": "Arial"}]},
    ])

shape_data = find_shape(slide12, "Google Shape;154;g3bdf68daae8_0_110")
if shape_data:
    update_textbox(shape_data, [
        {"runs": [{"text": "Подход к данным", "bold": True, "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Canonical SMILES → молекулярный граф (atom_num, chirality, bond_type, bond_dir)", "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Очистка/дедупликация RDKit, Balanced Scaffold Split (Bemis-Murcko)", "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "3 seeds → mean ± std, Adam (lr=1e-3), Early Stopping (patience=15), Batch 64", "size": Pt(13), "font": "Arial"}]},
    ])

shape_metrics = find_shape(slide12, "Google Shape;155;g3bdf68daae8_0_110")
if shape_metrics:
    update_textbox(shape_metrics, [
        {"runs": [{"text": "Метрики оценки", "bold": True, "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Качество: ROC-AUC, PR-AUC (mean по 12 задачам Tox21)", "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Обобщаемость: Tox21, BBBP, ClinTox — всё на scaffold split", "size": Pt(13), "font": "Arial"}]},
        {"runs": [{"text": "Объяснимость: AOPC@20% (faithfulness), IoU@20% (stability)", "size": Pt(13), "font": "Arial"}]},
    ])

# ============================================================
# SLIDE 13 (idx 12) — Источники Каболова: обновляем
# ============================================================
slide13 = slides[12]
shape_refs = find_shape(slide13, "Google Shape;162;g38c0d71ece1_0_356")
if shape_refs:
    refs = [
        "Hu, W. et al. (2020). Strategies for Pre-training GNN. ICLR 2020.",
        "Xu, K. et al. (2019). How Powerful are Graph Neural Networks? ICLR 2019.",
        "Gilmer, J. et al. (2017). Neural Message Passing for Quantum Chemistry. ICML 2017.",
        "Wu, Z. et al. (2018). MoleculeNet: A Benchmark for Molecular ML. Chemical Science.",
        "Ying, R. et al. (2019). GNNExplainer. NeurIPS 2019.",
        "Hu, W. et al. (2021). Open Graph Benchmark. NeurIPS 2020.",
        "Kendall, A. et al. (2018). Multi-Task Learning Using Uncertainty. CVPR 2018.",
        "Xiong, Z. et al. (2020). Pushing the Boundaries of Molecular Representation. J. Med. Chem.",
        "Lin, T.-Y. et al. (2017). Focal Loss for Dense Object Detection. ICCV 2017.",
        "Schütt, K. et al. (2017). SchNet. NeurIPS 2017.",
    ]
    pdata = [{"runs": [{"text": r, "size": Pt(14)}]} for r in refs]
    update_textbox(shape_refs, pdata)

# ============================================================
# Сохраняем
# ============================================================
prs.save(DST)
print(f"Сохранено: {DST}")
print("Обновлены слайды: 4 (решение), 5 (УЦП), 7 (прогресс), 9 (роль), 10 (гипотезы), 11 (анализ), 12 (методика), 13 (источники)")
